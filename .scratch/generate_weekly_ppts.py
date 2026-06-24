"""
Generate 3 weekly PPTs (Week 1/2/3) for 公务员考试学习跟踪系统.
Pure text layout, no images, academic defense style.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from lxml import etree
import os

# ---------------------------------------------------------------------------
# Theme constants
# ---------------------------------------------------------------------------
SLIDE_WIDTH = Inches(13.333)   # 16:9
SLIDE_HEIGHT = Inches(7.5)
HEADER_H = Inches(0.97)        # 70px ~ 0.97in
FOOTER_H = Inches(0.42)        # 30px
MARGIN_L = Inches(0.56)        # 40px
MARGIN_R = Inches(0.56)
MARGIN_T = HEADER_H + Inches(0.28)
MARGIN_B = SLIDE_HEIGHT - FOOTER_H - Inches(0.14)

COLORS = {
    'bg': RGBColor(0xFF, 0xFF, 0xFF),
    'header': RGBColor(0x00, 0x33, 0x66),
    'accent_red': RGBColor(0xCC, 0x00, 0x00),
    'accent_blue': RGBColor(0x00, 0x66, 0xCC),
    'text': RGBColor(0x33, 0x33, 0x33),
    'text_secondary': RGBColor(0x66, 0x66, 0x66),
    'text_tertiary': RGBColor(0x99, 0x99, 0x99),
    'card_bg': RGBColor(0xF5, 0xF7, 0xFA),
    'light_blue': RGBColor(0xE8, 0xF4, 0xFC),
}

FONT = 'Microsoft YaHei'


def set_slide_size(prs):
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT


def add_header(slide, title, section=''):
    """Add dark blue header bar with title."""
    # Header background
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, HEADER_H
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['header']
    header.line.fill.background()

    # Red accent bar
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.06), HEADER_H
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS['accent_red']
    accent.line.fill.background()

    # Section badge (top right)
    if section:
        badge_w = Inches(1.5)
        badge_h = Inches(0.35)
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            SLIDE_WIDTH - MARGIN_R - badge_w,
            (HEADER_H - badge_h) / 2,
            badge_w, badge_h
        )
        badge.fill.background()
        badge.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        badge.line.width = Pt(1)
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = section
        p.font.name = FONT
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)

    # Title
    title_box = slide.shapes.add_textbox(
        MARGIN_L + Inches(0.1), (HEADER_H - Inches(0.4)) / 2,
        SLIDE_WIDTH - MARGIN_L - MARGIN_R - Inches(1.7), Inches(0.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_footer(slide, page_num, total):
    """Add footer with page number."""
    # Footer line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_L, SLIDE_HEIGHT - FOOTER_H - Inches(0.08),
        SLIDE_WIDTH - MARGIN_L - MARGIN_R, Pt(1)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xD0, 0xD7, 0xE0)
    line.line.fill.background()

    # Page number
    box = slide.shapes.add_textbox(
        SLIDE_WIDTH - MARGIN_R - Inches(0.8), SLIDE_HEIGHT - FOOTER_H + Inches(0.02),
        Inches(0.8), Inches(0.3)
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = f'{page_num} / {total}'
    p.font.name = FONT
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['text_tertiary']
    p.alignment = PP_ALIGN.RIGHT


def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                 color=None, align=PP_ALIGN.LEFT, line_spacing=1.3):
    """Add a generic text box."""
    if color is None:
        color = COLORS['text']
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.line_spacing = line_spacing
    return box


def add_bullet_list(slide, left, top, width, height, items, font_size=18,
                    color=None, bullet_color=None):
    """Add bullet list."""
    if color is None:
        color = COLORS['text']
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f'•  {item}'
        p.font.name = FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(10)
        p.line_spacing = 1.3
        p.level = 0
    return box


def add_numbered_list(slide, left, top, width, height, items, font_size=18,
                      color=None):
    """Add numbered list."""
    if color is None:
        color = COLORS['text']
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f'{i+1}. {item}'
        p.font.name = FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(10)
        p.line_spacing = 1.3
    return box


def add_card(slide, left, top, width, height, title, items=None, font_size=16):
    """Add a light card with title and optional bullet items."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['card_bg']
    shape.line.color.rgb = RGBColor(0xD0, 0xD7, 0xE0)
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT
    p.font.size = Pt(font_size + 2)
    p.font.bold = True
    p.font.color.rgb = COLORS['header']
    p.space_after = Pt(8)

    if items:
        for item in items:
            p = tf.add_paragraph()
            p.text = f'• {item}'
            p.font.name = FONT
            p.font.size = Pt(font_size)
            p.font.color.rgb = COLORS['text']
            p.space_after = Pt(4)
            p.line_spacing = 1.25
    return shape


def add_kpi_cards(slide, kpis):
    """Add KPI cards in a row. kpis: list of (label, value, unit)."""
    n = len(kpis)
    gap = Inches(0.25)
    card_w = (SLIDE_WIDTH - MARGIN_L - MARGIN_R - gap * (n - 1)) / n
    card_h = Inches(1.6)
    top = MARGIN_T + Inches(0.4)

    for i, (label, value, unit) in enumerate(kpis):
        left = MARGIN_L + i * (card_w + gap)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS['light_blue']
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.clear()

        p = tf.paragraphs[0]
        p.text = value
        p.font.name = FONT
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = COLORS['accent_blue']
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(4)

        p = tf.add_paragraph()
        p.text = unit
        p.font.name = FONT
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['text_secondary']
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)

        p = tf.add_paragraph()
        p.text = label
        p.font.name = FONT
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['text']
        p.alignment = PP_ALIGN.CENTER


def slide_cover(slide, title, subtitle, info, section=''):
    """Cover slide."""
    # Full header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(3.2))
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['header']
    header.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.08), Inches(3.2))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS['accent_red']
    accent.line.fill.background()

    if section:
        add_text_box(slide, SLIDE_WIDTH - Inches(2.3), Inches(0.4), Inches(1.8), Inches(0.4),
                     section, font_size=16, bold=True,
                     color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

    add_text_box(slide, MARGIN_L, Inches(1.2), SLIDE_WIDTH - MARGIN_L - MARGIN_R, Inches(1.0),
                 title, font_size=48, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    add_text_box(slide, MARGIN_L, Inches(2.4), SLIDE_WIDTH - MARGIN_L - MARGIN_R, Inches(0.6),
                 subtitle, font_size=24, color=RGBColor(0xE8, 0xF4, 0xFC), align=PP_ALIGN.CENTER)

    # Decorative line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        (SLIDE_WIDTH - Inches(4.0)) / 2, Inches(3.5), Inches(4.0), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['accent_blue']
    line.line.fill.background()

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        (SLIDE_WIDTH - Inches(0.12)) / 2, Inches(3.5) - Inches(0.05), Inches(0.12), Inches(0.12))
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLORS['accent_blue']
    circle.line.fill.background()

    # Info block
    add_text_box(slide, MARGIN_L, Inches(4.2), SLIDE_WIDTH - MARGIN_L - MARGIN_R, Inches(2.0),
                 '\n'.join(info), font_size=20, color=COLORS['text_secondary'],
                 align=PP_ALIGN.CENTER, line_spacing=1.5)


def slide_toc(slide, items, section=''):
    """Table of contents slide."""
    add_header(slide, '汇报提纲', section)
    add_text_box(slide, MARGIN_L, MARGIN_T, Inches(2.5), Inches(0.6),
                 '目录', font_size=32, bold=True, color=COLORS['header'])

    content_top = MARGIN_T + Inches(0.9)
    content_h = MARGIN_B - content_top - Inches(0.3)
    box = slide.shapes.add_textbox(MARGIN_L, content_top,
        SLIDE_WIDTH - MARGIN_L - MARGIN_R, content_h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    for i, (num, title, desc) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f'0{num}    {title}'
        p.font.name = FONT
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = COLORS['header']
        p.space_after = Pt(4)

        p = tf.add_paragraph()
        p.text = f'        {desc}'
        p.font.name = FONT
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['text_secondary']
        p.space_after = Pt(16)


def slide_content(slide, title, bullets, section=''):
    """Standard content slide with title and bullets."""
    add_header(slide, title, section)
    add_bullet_list(slide, MARGIN_L, MARGIN_T + Inches(0.2),
        SLIDE_WIDTH - MARGIN_L - MARGIN_R, MARGIN_B - MARGIN_T - Inches(0.4),
        bullets, font_size=20)


def slide_two_column(slide, title, left_title, left_items, right_title, right_items, section=''):
    """Two-column content slide."""
    add_header(slide, title, section)
    col_w = (SLIDE_WIDTH - MARGIN_L - MARGIN_R - Inches(0.4)) / 2
    add_card(slide, MARGIN_L, MARGIN_T + Inches(0.2), col_w, Inches(0.45),
             left_title, font_size=18)
    add_bullet_list(slide, MARGIN_L, MARGIN_T + Inches(0.8), col_w,
        MARGIN_B - MARGIN_T - Inches(0.8), left_items, font_size=16)

    add_card(slide, MARGIN_L + col_w + Inches(0.4), MARGIN_T + Inches(0.2),
             col_w, Inches(0.45), right_title, font_size=18)
    add_bullet_list(slide, MARGIN_L + col_w + Inches(0.4), MARGIN_T + Inches(0.8),
        col_w, MARGIN_B - MARGIN_T - Inches(0.8), right_items, font_size=16)


def slide_kpi(slide, title, kpis, bullets, section=''):
    """Slide with KPI cards on top and bullets below."""
    add_header(slide, title, section)
    add_kpi_cards(slide, kpis)
    add_bullet_list(slide, MARGIN_L, MARGIN_T + Inches(2.2),
        SLIDE_WIDTH - MARGIN_L - MARGIN_R, MARGIN_B - MARGIN_T - Inches(2.4),
        bullets, font_size=18)


def slide_ending(slide, title, subtitle, info, section=''):
    """Thank you / Q&A slide."""
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['header']
    header.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.08), SLIDE_HEIGHT)
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS['accent_red']
    accent.line.fill.background()

    add_text_box(slide, MARGIN_L, Inches(2.4), SLIDE_WIDTH - MARGIN_L - MARGIN_R, Inches(1.0),
                 title, font_size=48, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, MARGIN_L, Inches(3.5), SLIDE_WIDTH - MARGIN_L - MARGIN_R, Inches(0.6),
                 subtitle, font_size=24, color=RGBColor(0xE8, 0xF4, 0xFC),
                 align=PP_ALIGN.CENTER)
    add_text_box(slide, MARGIN_L, Inches(4.5), SLIDE_WIDTH - MARGIN_L - MARGIN_R, Inches(1.5),
                 '\n'.join(info), font_size=16, color=RGBColor(0xCC, 0xCC, 0xCC),
                 align=PP_ALIGN.CENTER, line_spacing=1.5)


# ---------------------------------------------------------------------------
# Build Week 1 PPT
# ---------------------------------------------------------------------------
def build_week1():
    prs = Presentation()
    set_slide_size(prs)

    slides_data = []

    # Slide 1: Cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_cover(s, '公务员考试学习跟踪系统', 'Week 1 阶段汇报',
                ['答辩人：[答辩人姓名]', '指导老师：[指导老师]', '[学校 / 学院]', '2026年6月'])
    slides_data.append('cover')

    # Slide 2: TOC
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_toc(s, [
        (1, '项目背景与目标', '面向在校大学生的公务员备考跟踪系统'),
        (2, 'Week 1 概览', '搭建骨架、完成数据库与用户资源模块'),
        (3, '技术选型与架构', 'Flask + MySQL + 原生 HTML/CSS/JS'),
        (4, '数据库设计', '12 张核心表 + 11 个索引'),
        (5, '用户与资源模块', 'Session 认证 + 23 条备考资源'),
        (6, '后续计划', 'Week 2 将完成题库、计划、统计、推荐、答疑'),
    ], section='Week 1')
    slides_data.append('toc')

    # Slide 3: Background
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content(s, '项目背景与目标', [
        '目标用户：大一、大二备考公务员的在校大学生',
        '核心痛点：备考资料分散、学习进度难跟踪、弱项难定位',
        '项目目标：轻量级一站式公务员考试学习跟踪系统',
        '技术栈：Python Flask + MySQL 8.0 + 原生 HTML/CSS/JS',
        '开发周期：3 周，Week 1 聚焦基础设施与核心数据',
    ], section='Week 1')
    slides_data.append('content')

    # Slide 4: Week 1 Overview
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_kpi(s, 'Week 1 概览 — 从 0 到 1', [
        ('代码提交', '11', '次'),
        ('新增文件', '32', '个'),
        ('数据库表', '12', '张'),
        ('种子资源', '23', '条'),
    ], [
        'Day 1-2：项目初始化与 Flask 骨架搭建',
        'Day 3：12 张核心数据表设计与种子数据导入',
        'Day 4：用户注册 / 登录 / Session 认证',
        'Day 5：考试资源列表 / 详情 / 分类筛选',
        '当前整体进度约 25%：骨架已搭好，核心业务待 Week 2 实现',
    ], section='Week 1')
    slides_data.append('kpi')

    # Slide 5: Tech Stack
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_two_column(s, '技术选型与架构',
        '技术选型原则',
        [
            '简单、可演示、不引入新框架',
            'Flask：轻量、路由清晰、文档丰富',
            'MySQL 8.0：支持事务与外键',
            '原生 HTML/CSS/JS：降低答辩环境成本',
        ],
        '前后端分离架构',
        [
            '后端端口 5001，提供 RESTful API',
            '前端端口 8080，静态文件服务',
            '统一响应格式：{ success, data, message }',
            'CORS 配置支持本地双端口开发',
        ], section='Week 1')
    slides_data.append('two_col')

    # Slide 6: Database
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content(s, '数据库设计 — 12 张核心表', [
        '一次性定义 12 张表，覆盖用户、科目、资源、题目、答题、计划、进度、弱项、推荐、留言全链路',
        '核心表：users / subjects / resources / questions / answers / plans / plan_items / progress / weak_points / recommendations / comments / goals',
        '表间通过外键关联，如 answers 关联 users 与 questions',
        '建立 11 个索引，覆盖用户名、用户 ID、科目 ID、资源 ID、答题时间等高频查询字段',
        '种子数据：7 个科目、23 条资源、4 个演示账号，执行 init_db.py 一键重建',
    ], section='Week 1')
    slides_data.append('content')

    # Slide 7: User Module
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content(s, '用户模块 — 注册登录与 Session 认证', [
        '实现注册、登录、登出、获取当前用户 4 个 API',
        '密码使用 werkzeug 的 scrypt 哈希，不存明文',
        '登录成功后写入 Flask Session，浏览器 Cookie 自动携带 session_id',
        '前端登录页预填演示账号 root / 123456，方便快速进入系统',
        '封装 login_required / admin_required 装饰器，区分普通用户与管理员权限',
    ], section='Week 1')
    slides_data.append('content')

    # Slide 8: Resource Module
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content(s, '资源模块 — 23 条备考资源入库', [
        '资源分类：考试大纲、历年真题、模拟题、备考资料、政策公告',
        '每条资源关联科目，支持资源详情页查看与分类筛选',
        '前端资源库支持按类型标签切换、按科目筛选、关键词搜索',
        '管理员可上传新资源与批量删除',
        '资源详情页展示内容摘要，为 Week 2 题目与资源联动预留接口',
    ], section='Week 1')
    slides_data.append('content')

    # Slide 9: Next Steps + Ending
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content(s, '后续计划 — Week 2 重点工作', [
        '题库与练习：实现 200 道种子题的列表、筛选、答题与记录',
        '智能学习计划：根据备考周期、每日时长、科目权重生成每日任务',
        '学习进度统计：答题历史、正确率、完成率可视化',
        '弱项识别与推荐：基于答题历史自动定位弱项并推荐资源',
        '题目解析与答疑：每题提供解析，支持留言提问',
    ], section='Week 1')
    slides_data.append('content')

    # Slide 10: Thank you
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_ending(s, '感谢聆听', '欢迎老师提问与指导',
                 ['项目仓库：github.com/Aafff623/civil-service-exam-tracker',
                  '当前阶段：Week 1 完成，Week 2 进入核心业务模块开发'], section='Week 1')
    slides_data.append('ending')

    for i, _ in enumerate(slides_data, 1):
        add_footer(prs.slides[i-1], i, len(slides_data))

    return prs


# ---------------------------------------------------------------------------
# Build Week 2 PPT
# ---------------------------------------------------------------------------
def build_week2():
    prs = Presentation()
    set_slide_size(prs)

    slides_data = []

    # Slide 1: Cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_cover(s, '公务员考试学习跟踪系统', 'Week 2 阶段汇报',
                ['答辩人：[答辩人姓名]', '指导老师：[指导老师]', '[学校 / 学院]', '2026年6月'])
    slides_data.append('cover')

    # Slide 2: TOC
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_toc(s, [
        (1, 'Week 1 回顾', '基础设施、数据库、用户与资源模块'),
        (2, 'Week 2 概览', '完成 6 大核心业务模块'),
        (3, '题库与练习', '200 道题 + 单选/多选/判断'),
        (4, '智能学习计划', '科目权重 × 难度 + 弱项加权'),
        (5, '学习进度统计', '答题历史 + 正确率 + 可视化'),
        (6, '弱项识别与推荐', '答题量 ≥ 5 且正确率 < 60%'),
        (7, '解析与答疑', '每题解析 + 留言提问'),
        (8, '后续计划', 'Week 3 将完成 UI/UX、测试、文档整理'),
    ], section='Week 2')
    slides_data.append('toc')

    # Slide 3: Week 1 Review
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_kpi(s, 'Week 1 回顾 — 基础设施与核心数据', [
        ('提交', '11', '次'),
        ('文件', '32', '个'),
        ('数据表', '12', '张'),
        ('资源', '23', '条'),
    ], [
        '完成项目初始化、Flask 骨架、前后端分离架构',
        '设计 12 张核心数据表并导入种子数据',
        '实现用户注册 / 登录 / Session 认证',
        '实现考试资源列表 / 详情 / 分类筛选',
    ], section='Week 2')
    slides_data.append('kpi')

    # Slide 4: Week 2 Overview
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_kpi(s, 'Week 2 概览 — 核心业务逻辑', [
        ('种子题', '200', '道'),
        ('业务模块', '6', '个'),
        ('API 蓝图', '8', '个'),
        ('完成度', '≈75', '%'),
    ], [
        'Day 6：题库与练习模块',
        'Day 7：智能学习计划模块',
        'Day 8：学习进度统计模块',
        'Day 9：弱项识别与个性化推荐模块',
        'Day 10：题目解析与答疑模块',
    ], section='Week 2')
    slides_data.append('kpi')

    # Slide 5: Question Bank
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_two_column(s, '题库与练习模块',
        '题目数据',
        [
            '200 道种子题已入库',
            '单选 152 道 / 判断 34 道 / 多选 14 道',
            '覆盖 7 个科目',
            '通过 resource_id 与 23 条资源关联',
        ],
        '练习功能',
        [
            '支持科目 + 资料来源 + 题型三级筛选',
            '单选 / 多选 / 判断三种答题交互',
            '多选答案规范化比对，避免 AC/CA 判错',
            '答题后即时展示结果与解析',
        ], section='Week 2')
    slides_data.append('two_col')

    # Slide 6: Study Plan
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content(s, '智能学习计划模块', [
        '输入：备考周期、每日可用时长、科目权重',
        '算法：优先级分数 = 科目权重 × 难度系数 + 弱项加权 50%',
        '输出：每日任务卡片，自动划分三阶段（基础 30% / 强化 40% / 冲刺 30%）',
        '每日任务数：max(2, min(5, daily_minutes // 30))',
        '重新生成计划时事务内先删后插，保证数据一致性',
    ], section='Week 2')
    slides_data.append('content')

    # Slide 7: Progress Statistics
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content(s, '学习进度统计模块', [
        '记录每次答题历史，实时计算正确率与完成率',
        'Dashboard 接入今日任务、弱项提醒、考试倒计时',
        '统计页提供 KPI 看板、学习时长折线图、正确率趋势',
        '打卡日历展示每日学习记录，支持连续打卡统计',
        '新用户无数据时展示 Demo 数据兜底，页面不空',
    ], section='Week 2')
    slides_data.append('content')

    # Slide 8: Weakness & Recommendation
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_two_column(s, '弱项识别与个性化推荐',
        '弱项识别规则',
        [
            '答题量 ≥ 5 且正确率 < 60%',
            '每次答题后实时更新 weak_points 表',
            '正确率动态计算，弱项状态自动变化',
        ],
        '推荐逻辑',
        [
            '基于弱项科目匹配相关备考资源',
            '推荐同类练习题巩固薄弱点',
            '新用户冷启动时回退到通用入门推荐',
        ], section='Week 2')
    slides_data.append('two_col')

    # Slide 9: QA Module
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content(s, '解析与答疑模块', [
        '每道题均提供答案解析与答题技巧',
        '支持学生针对题目提交留言提问',
        'Dashboard 接入评论提交入口',
        '答疑回复当前为规则回复，满足课程设计可控范围',
        '留言数据与题目、用户关联，便于后续追踪',
    ], section='Week 2')
    slides_data.append('content')

    # Slide 10: Current Status
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_kpi(s, '当前完成度', [
        ('核心模块', '7/7', '个'),
        ('数据库表', '12', '张'),
        ('种子题库', '200', '题'),
        ('备考资源', '23', '条'),
    ], [
        '全部 7 个核心模块已实现并跑通',
        '功能逻辑完成，但 UI/UX 尚未统一打磨',
        '测试覆盖与文档整理将在 Week 3 进行',
    ], section='Week 2')
    slides_data.append('kpi')

    # Slide 11: Next Steps
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content(s, '后续计划 — Week 3 重点工作', [
        'UI/UX 统一：迁移 Aceternity 风格 surface、优化页面过渡与动画',
        '模块 Review：Dashboard、统计页、资源详情、个人资料页细节打磨',
        '测试与调试：PowerShell API 冒烟测试、修复时区与加载时序问题',
        '演示准备：题库与资源关联、多选支持完善、换机部署 README',
        '文档归档：README.md、PROJECT_GUIDE.md、DEVELOPMENT_TIMELINE.md',
    ], section='Week 2')
    slides_data.append('content')

    # Slide 12: Thank you
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_ending(s, '感谢聆听', '欢迎老师提问与指导',
                 ['项目仓库：github.com/Aafff623/civil-service-exam-tracker',
                  '当前阶段：Week 2 完成核心业务，Week 3 进入完善与收尾'], section='Week 2')
    slides_data.append('ending')

    for i, _ in enumerate(slides_data, 1):
        add_footer(prs.slides[i-1], i, len(slides_data))

    return prs


# ---------------------------------------------------------------------------
# Build Week 3 PPT
# ---------------------------------------------------------------------------
def build_week3():
    prs = Presentation()
    set_slide_size(prs)

    slides_data = []

    # Slide 1: Cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_cover(s, '公务员考试学习跟踪系统', '课程设计答辩汇报',
                ['答辩人：[答辩人姓名]', '指导老师：[指导老师]', '[学校 / 学院]', '2026年6月'])
    slides_data.append('cover')

    # Slide 2: TOC
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_toc(s, [
        (1, '项目背景与目标', '面向在校大学生的公务员备考跟踪系统'),
        (2, 'Week 1 回顾', '基础设施、数据库、用户与资源模块'),
        (3, 'Week 2 回顾', '题库、计划、统计、推荐、答疑'),
        (4, 'Week 3 概览', 'UI/UX、测试、文档、最终成果'),
        (5, 'UI/UX 统一与动画', 'Aceternity 风格与交互体验'),
        (6, '测试与演示数据', 'API 冒烟测试与 200 题题库'),
        (7, '文档整理与归档', 'README、PROJECT_GUIDE、时间线'),
        (8, '最终成果与展望', '7 大模块 + 可运行系统'),
    ], section='Week 3')
    slides_data.append('toc')

    # Slide 3: Background
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content(s, '项目背景与目标', [
        '目标用户：大一、大二备考公务员的在校大学生',
        '核心痛点：备考资料分散、学习进度难跟踪、弱项难定位',
        '项目目标：轻量级一站式公务员考试学习跟踪系统',
        '技术栈：Python Flask + MySQL 8.0 + 原生 HTML/CSS/JS',
        '开发周期：3 周，从骨架搭建到可答辩系统',
    ], section='Week 3')
    slides_data.append('content')

    # Slide 4: Week 1 Review
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_kpi(s, 'Week 1 回顾 — 基础设施', [
        ('提交', '11', '次'),
        ('表', '12', '张'),
        ('资源', '23', '条'),
        ('蓝图', '2', '个'),
    ], [
        '完成项目初始化、Flask 骨架、前后端分离架构',
        '设计 12 张核心数据表并导入种子数据',
        '实现用户注册 / 登录 / Session 认证',
        '实现考试资源列表 / 详情 / 分类筛选',
    ], section='Week 3')
    slides_data.append('kpi')

    # Slide 5: Week 2 Review
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_kpi(s, 'Week 2 回顾 — 核心业务', [
        ('题库', '200', '题'),
        ('模块', '6', '个'),
        ('蓝图', '10', '个'),
        ('完成度', '≈75', '%'),
    ], [
        '完成题库与练习、智能学习计划、学习进度统计',
        '完成弱项识别与个性化推荐、解析与答疑',
        'Dashboard 接入今日任务、弱项提醒、评论提交',
        '全部 7 个核心模块功能逻辑已实现',
    ], section='Week 3')
    slides_data.append('kpi')

    # Slide 6: Week 3 Overview
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_kpi(s, 'Week 3 概览 — 完善、Review、演示', [
        ('提交', '42', '次'),
        ('文件', '109', '个'),
        ('测试端点', '18+', '个'),
        ('完成度', '100', '%'),
    ], [
        'Day 11：UI/UX 统一，Aceternity 风格 surface 迁移',
        'Day 12：模块 Review，Dashboard 与统计页打磨',
        'Day 13：测试与调试，PowerShell API 冒烟测试',
        'Day 14：演示准备，题库扩容至 200 题、资源关联',
        'Day 15：收尾，README polish、CLAUDE.md、.cursorrules',
    ], section='Week 3')
    slides_data.append('kpi')

    # Slide 7: UI/UX
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content(s, 'UI/UX 统一与动画交互', [
        '设计系统：从 GPT 生成资产迁移到 Aceternity 风格 surfaces.css',
        '加载 veil、页面过渡、GSAP KPI 数字滚动动画',
        'Dashboard 今日任务勾选 FLIP 动画',
        '学习时长热力图支持 5 级色阶与悬停 tooltip',
        '860px 断点下 sidebar 变为抽屉式导航，支持移动端浏览',
    ], section='Week 3')
    slides_data.append('content')

    # Slide 8: Testing
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content(s, '测试与演示数据准备', [
        '编写 PowerShell 冒烟测试脚本，覆盖 18+ API 端点',
        '修复模块加载时序、本地时区过滤、多选答案比对等问题',
        '题库从初始版本扩容至 200 题（单选 152 / 判断 34 / 多选 14）',
        '题目与资源通过 resource_id 关联，支持资源详情页直接练习',
        '完善多选题型交互与答案规范化比对',
    ], section='Week 3')
    slides_data.append('content')

    # Slide 9: Documentation
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content(s, '文档整理与项目归档', [
        'README.md：环境准备、开机运行、换机实测、维护规范',
        'docs/PROJECT_GUIDE.md：详细项目使用指南',
        'docs/DEVELOPMENT_TIMELINE.md：三周开发进度与关键节点',
        '.cursorrules：Cursor 项目上下文与代码规范',
        '所有文档与 init_db.sql 同步更新，确保换机可复现',
    ], section='Week 3')
    slides_data.append('content')

    # Slide 10: Final Results
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_kpi(s, '最终成果', [
        ('模块', '7', '个'),
        ('数据表', '12', '张'),
        ('题库', '200', '题'),
        ('资源', '23', '条'),
    ], [
        '7 个核心模块全部实现并联调通过',
        '10 个后端 Blueprint，后端 5001 + 前端 8080',
        '200 题 + 23 资源 + 4 个演示账号',
        '约 50+ 次 commit，系统可运行、可演示、可换机部署',
    ], section='Week 3')
    slides_data.append('kpi')

    # Slide 11: Core Data
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_two_column(s, '核心数据一览',
        '代码与模块',
        [
            '后端 Blueprint：10 个',
            '前端页面：11 个 HTML',
            'API 端点：18+ 个',
            'Git commit：50+ 次',
        ],
        '数据与内容',
        [
            '数据库表：12 张',
            '种子题库：200 道题',
            '备考资源：23 条',
            '演示账号：4 个',
        ], section='Week 3')
    slides_data.append('two_col')

    # Slide 12: Limitations
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content(s, '后续计划与不足', [
        '考试时间线：当前为静态文案，可扩展为接入考试公告 API',
        'AI 答疑：当前为规则回复，可接入 LLM 提供更智能解答',
        '移动端体验：已做响应式适配，但部分页面仍可继续优化',
        '性能优化：大数据量下可引入缓存与分页优化',
        '商业化扩展：可增加学习小组、排行榜、错题本等功能',
    ], section='Week 3')
    slides_data.append('content')

    # Slide 13: Demo Info
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content(s, '系统演示信息', [
        '启动方式：后端 python app.py（端口 5001）+ 前端 python -m http.server 8080',
        '访问地址：http://localhost:8080/login.html',
        '演示账号：root / 123456',
        '主要页面：登录页、Dashboard、资源库、题库练习、学习计划、统计页、推荐页、个人资料页',
        '验收路径：登录 → 浏览资源 → 开始练习 → 查看统计 → 生成推荐',
    ], section='Week 3')
    slides_data.append('content')

    # Slide 14: Thank you
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_ending(s, '感谢聆听', '欢迎老师提问与指导',
                 ['项目仓库：github.com/Aafff623/civil-service-exam-tracker',
                  '演示地址：http://localhost:8080/login.html',
                  '账号：root / 123456'], section='Week 3')
    slides_data.append('ending')

    for i, _ in enumerate(slides_data, 1):
        add_footer(prs.slides[i-1], i, len(slides_data))

    return prs


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
if __name__ == '__main__':
    out_dir = r'D:\OneDrive\Desktop\课设接单\civil-service-exam-tracker\docs\ppt'
    os.makedirs(out_dir, exist_ok=True)

    ppts = [
        ('Week1_公务员考试学习跟踪系统_阶段汇报.pptx', build_week1),
        ('Week2_公务员考试学习跟踪系统_阶段汇报.pptx', build_week2),
        ('Week3_公务员考试学习跟踪系统_答辩.pptx', build_week3),
    ]

    for filename, builder in ppts:
        path = os.path.join(out_dir, filename)
        prs = builder()
        prs.save(path)
        print(f'Saved: {path}')

    print('\nAll 3 PPTs generated successfully.')
