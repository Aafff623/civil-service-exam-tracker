"""
Generate final overview defense PPT for 公务员考试学习跟踪系统.
18 slides, text + image placeholders, academic defense style.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ---------------------------------------------------------------------------
# Theme constants (compact version)
# ---------------------------------------------------------------------------
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
HEADER_H = Inches(0.72)
FOOTER_H = Inches(0.28)
MARGIN_L = Inches(0.5)
MARGIN_R = Inches(0.5)
MARGIN_T = HEADER_H + Inches(0.14)
MARGIN_B = SLIDE_HEIGHT - FOOTER_H - Inches(0.08)

COLORS = {
    'bg': RGBColor(0xFF, 0xFF, 0xFF),
    'header': RGBColor(0x00, 0x33, 0x66),
    'accent_red': RGBColor(0xCC, 0x00, 0x00),
    'accent_blue': RGBColor(0x00, 0x66, 0xCC),
    'text': RGBColor(0x33, 0x33, 0x33),
    'text_secondary': RGBColor(0x66, 0x66, 0x66),
    'text_tertiary': RGBColor(0x99, 0x99, 0x99),
    'card_bg': RGBColor(0xF5, 0xF7, 0xFA),
    'light_blue': RGBColor(0xE8, 0xF4, 0xFC),
    'placeholder_bg': RGBColor(0xF0, 0xF7, 0xFF),
}

FONT = 'Microsoft YaHei'


def set_slide_size(prs):
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT


def add_header(slide, title, section=''):
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, HEADER_H)
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['header']
    header.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.06), HEADER_H)
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS['accent_red']
    accent.line.fill.background()

    if section:
        badge_w = Inches(1.3)
        badge_h = Inches(0.32)
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            SLIDE_WIDTH - MARGIN_R - badge_w, (HEADER_H - badge_h) / 2,
            badge_w, badge_h
        )
        badge.fill.background()
        badge.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        badge.line.width = Pt(1)
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = section
        p.font.name = FONT
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER

    title_box = slide.shapes.add_textbox(
        MARGIN_L + Inches(0.08), (HEADER_H - Inches(0.36)) / 2,
        SLIDE_WIDTH - MARGIN_L - MARGIN_R - Inches(1.5), Inches(0.45)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_footer(slide, page_num, total):
    footer_y = SLIDE_HEIGHT - FOOTER_H
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_L, footer_y + Inches(0.04),
        SLIDE_WIDTH - MARGIN_L - MARGIN_R, Pt(1)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xD0, 0xD7, 0xE0)
    line.line.fill.background()

    box = slide.shapes.add_textbox(
        SLIDE_WIDTH - MARGIN_R - Inches(0.7), footer_y + Inches(0.06),
        Inches(0.7), Inches(0.18)
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = f'{page_num} / {total}'
    p.font.name = FONT
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS['text_tertiary']
    p.alignment = PP_ALIGN.RIGHT


def add_text_box(slide, left, top, width, height, text, font_size=15, bold=False,
                 color=None, align=PP_ALIGN.LEFT, line_spacing=1.1):
    if color is None:
        color = COLORS['text']
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.line_spacing = line_spacing
    return box


def add_bullet_list(slide, left, top, width, height, items, font_size=15, color=None):
    if color is None:
        color = COLORS['text']
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f'•  {item}'
        p.font.name = FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(3)
        p.line_spacing = 1.10
        p.level = 0
    return box


def add_numbered_list(slide, left, top, width, height, items, font_size=15, color=None):
    if color is None:
        color = COLORS['text']
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f'{i+1}. {item}'
        p.font.name = FONT
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(3)
        p.line_spacing = 1.10
    return box


def add_kpi_cards(slide, kpis, top_offset=Inches(0.15)):
    n = len(kpis)
    gap = Inches(0.18)
    card_w = (SLIDE_WIDTH - MARGIN_L - MARGIN_R - gap * (n - 1)) / n
    card_h = Inches(1.05)
    top = MARGIN_T + top_offset

    for i, (label, value, unit) in enumerate(kpis):
        left = MARGIN_L + i * (card_w + gap)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS['light_blue']
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        tf.clear()

        p = tf.paragraphs[0]
        p.text = value
        p.font.name = FONT
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = COLORS['accent_blue']
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(1)

        p = tf.add_paragraph()
        p.text = unit
        p.font.name = FONT
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['text_secondary']
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(2)

        p = tf.add_paragraph()
        p.text = label
        p.font.name = FONT
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['text']
        p.alignment = PP_ALIGN.CENTER


def add_image_placeholder(slide, left, top, width, height, label):
    """Add a placeholder box for screenshot/diagram."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['placeholder_bg']
    shape.line.color.rgb = COLORS['accent_blue']
    shape.line.width = Pt(1.5)

    # Dashed line effect via secondary transparent overlay not supported simply;
    # use label to indicate placeholder.
    add_text_box(slide, left, top + height / 2 - Inches(0.12),
                 width, Inches(0.24),
                 f'[ {label} ]', font_size=12, bold=True,
                 color=COLORS['accent_blue'], align=PP_ALIGN.CENTER)
    return shape


def slide_cover(slide, title, subtitle, info):
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(3.0))
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['header']
    header.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.08), Inches(3.0))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS['accent_red']
    accent.line.fill.background()

    add_text_box(slide, MARGIN_L, Inches(1.0), SLIDE_WIDTH - MARGIN_L - MARGIN_R, Inches(0.9),
                 title, font_size=42, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    add_text_box(slide, MARGIN_L, Inches(2.1), SLIDE_WIDTH - MARGIN_L - MARGIN_R, Inches(0.5),
                 subtitle, font_size=20, color=RGBColor(0xE8, 0xF4, 0xFC), align=PP_ALIGN.CENTER)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        (SLIDE_WIDTH - Inches(4.0)) / 2, Inches(3.1), Inches(4.0), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['accent_blue']
    line.line.fill.background()

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        (SLIDE_WIDTH - Inches(0.12)) / 2, Inches(3.1) - Inches(0.05), Inches(0.12), Inches(0.12))
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLORS['accent_blue']
    circle.line.fill.background()

    add_text_box(slide, MARGIN_L, Inches(3.6), SLIDE_WIDTH - MARGIN_L - MARGIN_R, Inches(2.0),
                 '\n'.join(info), font_size=16, color=COLORS['text_secondary'],
                 align=PP_ALIGN.CENTER, line_spacing=1.25)


def slide_toc(slide, items):
    add_header(slide, '汇报提纲')
    add_text_box(slide, MARGIN_L, MARGIN_T, Inches(2.0), Inches(0.45),
                 '目录', font_size=24, bold=True, color=COLORS['header'])

    content_top = MARGIN_T + Inches(0.6)
    content_h = MARGIN_B - content_top - Inches(0.1)
    col_w = (SLIDE_WIDTH - MARGIN_L - MARGIN_R - Inches(0.3)) / 2

    mid = (len(items) + 1) // 2
    left_items = items[:mid]
    right_items = items[mid:]

    for col_idx, col_items in enumerate([left_items, right_items]):
        left = MARGIN_L + col_idx * (col_w + Inches(0.3))
        box = slide.shapes.add_textbox(left, content_top, col_w, content_h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()

        for i, (num, title, desc) in enumerate(col_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f'{num:02d}  {title}'
            p.font.name = FONT
            p.font.size = Pt(15)
            p.font.bold = True
            p.font.color.rgb = COLORS['header']
            p.space_after = Pt(1)

            p = tf.add_paragraph()
            p.text = f'     {desc}'
            p.font.name = FONT
            p.font.size = Pt(10)
            p.font.color.rgb = COLORS['text_secondary']
            p.space_after = Pt(8)


def slide_content(slide, title, bullets, section=''):
    add_header(slide, title, section)
    add_bullet_list(slide, MARGIN_L, MARGIN_T + Inches(0.12),
        SLIDE_WIDTH - MARGIN_L - MARGIN_R, MARGIN_B - MARGIN_T - Inches(0.25),
        bullets, font_size=15)


def slide_content_with_image(slide, title, bullets, placeholder_label, section=''):
    """Left bullets (55%) + right image placeholder (45%)."""
    add_header(slide, title, section)

    left_w = Inches(6.7)
    right_w = SLIDE_WIDTH - MARGIN_L - MARGIN_R - left_w - Inches(0.25)
    content_top = MARGIN_T + Inches(0.12)
    content_h = MARGIN_B - MARGIN_T - Inches(0.25)

    add_bullet_list(slide, MARGIN_L, content_top, left_w, content_h,
                    bullets, font_size=14)

    ph_left = MARGIN_L + left_w + Inches(0.25)
    ph_top = content_top + Inches(0.2)
    ph_h = min(content_h - Inches(0.4), Inches(4.2))
    add_image_placeholder(slide, ph_left, ph_top, right_w, ph_h, placeholder_label)


def slide_kpi(slide, title, kpis, bullets, section=''):
    add_header(slide, title, section)
    add_kpi_cards(slide, kpis)
    add_bullet_list(slide, MARGIN_L, MARGIN_T + Inches(1.4),
        SLIDE_WIDTH - MARGIN_L - MARGIN_R, MARGIN_B - MARGIN_T - Inches(1.55),
        bullets, font_size=15)


def slide_two_column(slide, title, left_title, left_items, right_title, right_items, section=''):
    add_header(slide, title, section)
    col_w = (SLIDE_WIDTH - MARGIN_L - MARGIN_R - Inches(0.35)) / 2

    # Column headers
    add_text_box(slide, MARGIN_L, MARGIN_T + Inches(0.08), col_w, Inches(0.32),
                 left_title, font_size=16, bold=True, color=COLORS['header'])
    add_bullet_list(slide, MARGIN_L, MARGIN_T + Inches(0.45), col_w,
        MARGIN_B - MARGIN_T - Inches(0.55), left_items, font_size=14)

    add_text_box(slide, MARGIN_L + col_w + Inches(0.35), MARGIN_T + Inches(0.08), col_w, Inches(0.32),
                 right_title, font_size=16, bold=True, color=COLORS['header'])
    add_bullet_list(slide, MARGIN_L + col_w + Inches(0.35), MARGIN_T + Inches(0.45), col_w,
        MARGIN_B - MARGIN_T - Inches(0.55), right_items, font_size=14)


def slide_ending(slide, title, subtitle, info):
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS['header']
    header.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.08), SLIDE_HEIGHT)
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS['accent_red']
    accent.line.fill.background()

    add_text_box(slide, MARGIN_L, Inches(2.3), SLIDE_WIDTH - MARGIN_L - MARGIN_R, Inches(0.9),
                 title, font_size=42, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    add_text_box(slide, MARGIN_L, Inches(3.2), SLIDE_WIDTH - MARGIN_L - MARGIN_R, Inches(0.5),
                 subtitle, font_size=20, color=RGBColor(0xE8, 0xF4, 0xFC), align=PP_ALIGN.CENTER)
    add_text_box(slide, MARGIN_L, Inches(4.0), SLIDE_WIDTH - MARGIN_L - MARGIN_R, Inches(1.6),
                 '\n'.join(info), font_size=14, color=RGBColor(0xCC, 0xCC, 0xCC),
                 align=PP_ALIGN.CENTER, line_spacing=1.2)


# ---------------------------------------------------------------------------
# Build final overview PPT
# ---------------------------------------------------------------------------
def build_final():
    prs = Presentation()
    set_slide_size(prs)
    slides = []

    # Slide 1: Cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_cover(s, '公务员考试学习跟踪系统', '课程设计答辩汇报',
                ['答辩人：[答辩人姓名]', '指导老师：[指导老师]', '[学校 / 学院]', '2026年6月'])
    slides.append('cover')

    # Slide 2: TOC
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_toc(s, [
        (1, '项目背景与目标', '面向在校大学生的公务员备考跟踪系统'),
        (2, '技术选型与架构', 'Flask + MySQL + 原生 HTML/CSS/JS'),
        (3, '数据库设计', '12 张核心表 + 11 个索引'),
        (4, '用户与资源模块', 'Session 认证 + 23 条备考资源'),
        (5, '题库与练习', '200 道题 + 单选/多选/判断'),
        (6, '智能学习计划', '权重 × 难度 + 弱项加权'),
        (7, '进度统计与可视化', '答题历史 + 正确率 + 图表'),
        (8, '弱项识别与推荐', '答题量 ≥ 5 且正确率 < 60%'),
        (9, '解析与答疑', '每题解析 + 留言提问'),
        (10, 'UI/UX 与测试', 'Aceternity 风格 + API 冒烟测试'),
        (11, '最终成果', '7 大模块 + 可运行系统'),
        (12, '致谢 / Q&A', '感谢聆听，欢迎提问'),
    ])
    slides.append('toc')

    # Slide 3: Background
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content(s, '项目背景与目标', [
        '目标用户：大一、大二备考公务员的在校大学生',
        '核心痛点：备考资料分散、学习进度难跟踪、弱项难定位',
        '项目目标：轻量级一站式公务员考试学习跟踪系统',
        '技术栈：Python Flask + MySQL 8.0 + 原生 HTML/CSS/JS',
        '开发周期：3 周，从骨架搭建到可答辩系统',
    ])
    slides.append('content')

    # Slide 4: Tech Stack + Architecture diagram placeholder
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content_with_image(s, '技术选型与架构',
        [
            '后端：Flask 3.0，RESTful API，Session/Cookie 认证',
            '数据库：MySQL 8.0+，12 张表覆盖 7 大模块',
            '前端：原生 HTML/CSS/JS，无额外框架',
            '架构：前后端分离，后端 5001 + 前端 8080',
            '部署：init_db.py 一键导入种子数据，换机方便',
        ], '此处插入：系统架构图')
    slides.append('content_img')

    # Slide 5: Database + ER diagram placeholder
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content_with_image(s, '数据库设计',
        [
            '12 张核心表：users / subjects / resources / questions / answers',
            'plans / plan_items / progress / weak_points / recommendations / comments / goals',
            '外键关联：answers 关联 users 与 questions',
            '11 个索引覆盖高频查询字段',
            '种子数据：7 科目、23 资源、4 演示账号',
        ], '此处插入：数据库 ER 图')
    slides.append('content_img')

    # Slide 6: User & Resource + screenshots placeholder
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content_with_image(s, '用户与资源模块',
        [
            '用户：注册 / 登录 / Session 认证 / scrypt 密码哈希',
            '权限：login_required 与 admin_required 装饰器',
            '资源：23 条备考资源，分 5 大类型',
            '资源库：类型标签 + 科目筛选 + 关键词搜索',
            '演示账号：root / 123456',
        ], '此处插入：登录页 / 资源库截图')
    slides.append('content_img')

    # Slide 7: Week 1 Summary
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content(s, 'Week 1 小结', [
        '完成项目骨架、数据库设计、用户与资源模块',
        '产出 11 次提交、32 个新增文件、12 张表、23 条资源',
        '为 Week 2 业务模块预留完整外键与命名规范',
        '技术选型与前后端分离架构确定',
    ])
    slides.append('content')

    # Slide 8: Question Bank + screenshot placeholder
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content_with_image(s, '题库与练习模块',
        [
            '200 道种子题：单选 152 / 判断 34 / 多选 14',
            '覆盖 7 个科目，通过 resource_id 与资源关联',
            '三级筛选：科目 + 资料来源 + 题型',
            '多选答案规范化比对，避免 AC/CA 判错',
            '答题后即时展示结果与解析',
        ], '此处插入：题库练习页截图')
    slides.append('content_img')

    # Slide 9: Study Plan + screenshot placeholder
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content_with_image(s, '智能学习计划',
        [
            '输入：备考周期、每日可用时长、科目权重',
            '算法：优先级 = 科目权重 × 难度 + 弱项加权 50%',
            '输出：每日任务卡片，三阶段划分',
            '基础 30% / 强化 40% / 冲刺 30%',
            '重新生成时事务内先删后插',
        ], '此处插入：学习计划页截图')
    slides.append('content_img')

    # Slide 10: Progress + screenshot placeholder
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content_with_image(s, '进度统计与可视化',
        [
            '记录答题历史，实时计算正确率与完成率',
            'Dashboard 接入今日任务与进度概览',
            '统计页：KPI 看板 + 学习时长折线图',
            '正确率趋势 + 打卡日历 + 连续打卡',
            '新用户无数据时 Demo 数据兜底',
        ], '此处插入：统计页截图')
    slides.append('content_img')

    # Slide 11: Weakness + Recommendation + screenshot placeholder
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content_with_image(s, '弱项识别与个性化推荐',
        [
            '弱项规则：答题量 ≥ 5 且正确率 < 60%',
            '每次答题后实时更新 weak_points 表',
            '自动推荐相关资源与同类练习题',
            '推荐页基于答题历史动态生成',
            '新用户冷启动回退到通用入门推荐',
        ], '此处插入：推荐页截图')
    slides.append('content_img')

    # Slide 12: QA + screenshot placeholder
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content_with_image(s, '解析与答疑模块',
        [
            '每题提供答案解析与答题技巧',
            '支持学生针对题目提交留言提问',
            'Dashboard 接入评论提交入口',
            'AI 答疑为规则回复，课程设计可控',
            '留言与题目、用户关联，便于追踪',
        ], '此处插入：答疑页截图')
    slides.append('content_img')

    # Slide 13: Week 2 Summary
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content(s, 'Week 2 小结', [
        '完成 6 大核心业务模块：题库、计划、统计、推荐、答疑',
        '200 道题入库，Dashboard 接入核心业务数据',
        '弱项识别与推荐逻辑跑通',
        '为 Week 3 UI/UX 统一与测试打下基础',
    ])
    slides.append('content')

    # Slide 14: UI/UX + Dashboard screenshot placeholder
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content_with_image(s, 'UI/UX 统一与动画交互',
        [
            'Aceternity 风格 surface 迁移',
            '加载 veil、页面过渡、GSAP 数字动画',
            'Dashboard 今日任务勾选 FLIP 动画',
            '学习时长热力图支持 5 级色阶',
            '860px 断点下 sidebar 变抽屉式导航',
        ], '此处插入：Dashboard 截图')
    slides.append('content_img')

    # Slide 15: Testing
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content(s, '测试、Review 与演示数据', [
        'PowerShell 冒烟测试覆盖 18+ API 端点',
        '修复模块加载时序、本地时区过滤、多选比对等问题',
        '题库从初始版扩容至 200 题',
        '题目与资源通过 resource_id 关联',
        '多选题型交互与答案规范化完善',
    ])
    slides.append('content')

    # Slide 16: Documentation
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_content(s, '文档整理与项目归档', [
        'README.md：环境准备 / 开机运行 / 换机实测 / 维护规范',
        'docs/PROJECT_GUIDE.md：详细使用指南',
        'docs/DEVELOPMENT_TIMELINE.md：三周开发进度',
        '.cursorrules：Cursor 项目上下文',
        '所有文档与 init_db.sql 同步更新',
    ])
    slides.append('content')

    # Slide 17: Final Results + multiple screenshot placeholders
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_kpi(s, '最终成果与核心数据', [
        ('核心模块', '7', '个'),
        ('数据表', '12', '张'),
        ('题库', '200', '题'),
        ('资源', '23', '条'),
    ], [
        '7 个核心模块全部实现并联调通过',
        '10 个后端 Blueprint，后端 5001 + 前端 8080',
        '约 50+ 次 commit，可运行、可演示、可换机部署',
        '演示账号：root / 123456',
    ])

    # Add screenshot placeholders below KPI cards
    ph_y = MARGIN_T + Inches(2.6)
    ph_w = (SLIDE_WIDTH - MARGIN_L - MARGIN_R - Inches(0.6)) / 4
    ph_h = Inches(1.55)
    labels = ['登录页', 'Dashboard', '资源库', '题库页']
    for i, label in enumerate(labels):
        ph_x = MARGIN_L + i * (ph_w + Inches(0.2))
        add_image_placeholder(s, ph_x, ph_y, ph_w, ph_h, f'此处插入：{label}截图')
    slides.append('kpi_img')

    # Slide 18: Thank you
    s = prs.slides.add_slide(prs.slide_layouts[6])
    slide_ending(s, '感谢聆听', '欢迎老师提问与指导',
                 ['项目仓库：github.com/Aafff623/civil-service-exam-tracker',
                  '演示地址：http://localhost:8080/login.html',
                  '账号：root / 123456'])
    slides.append('ending')

    for i, _ in enumerate(slides, 1):
        add_footer(prs.slides[i-1], i, len(slides))

    return prs


if __name__ == '__main__':
    out_dir = r'D:\OneDrive\Desktop\课设接单\civil-service-exam-tracker\docs\ppt'
    os.makedirs(out_dir, exist_ok=True)
    path = os.path.join(out_dir, '公务员考试学习跟踪系统答辩.pptx')

    # Remove old file if exists to avoid permission issues
    if os.path.exists(path):
        try:
            os.remove(path)
        except PermissionError:
            print(f'Warning: cannot remove existing file {path}, it may be open.')
            path = os.path.join(out_dir, '公务员考试学习跟踪系统答辩_新版.pptx')

    prs = build_final()
    prs.save(path)
    print(f'Saved: {path}')
